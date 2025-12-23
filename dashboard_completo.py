# dashboard_completo.py
import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
import numpy as np
import requests
from PIL import Image
import io
import os
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

# =============================================
# CONFIGURACIÓN STREAMLIT
# =============================================
st.set_page_config(
    page_title="Dashboard de Producción - Adimatec",
    page_icon="🏭",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Limpiar cache
st.cache_data.clear()

# Cargar logo
@st.cache_data
def load_logo(url):
    try:
        response = requests.get(url, timeout=10)
        image = Image.open(io.BytesIO(response.content))
        return image
    except:
        return None

logo = load_logo("https://i.postimg.cc/hjfVhfXf/Logo-Adimatec.jpg")

# Título principal con logo
col_logo, col_title, col_icon = st.columns([1, 3, 1])
with col_logo:
    if logo:
        st.image(logo, width=100)
with col_title:
    st.title("Dashboard de Producción - Adimatec")
with col_icon:
    st.markdown(
        """
        <div style='text-align: right; margin-top: 20px;'>
            <span style='font-size: 2em;'>🏭</span>
        </div>
        """,
        unsafe_allow_html=True
    )

st.markdown("---")

@st.cache_data(ttl=300, show_spinner="Cargando datos desde Google Sheets...")
def load_data():
    """Cargar datos desde Google Sheets"""
    try:
        # Sheet ID (la parte larga después de /d/)
        sheet_id = "17eEYewfzoBZXkFWBm5DOJp3IuvHg9WvN"
        
        # URLs CORREGIDAS - formato de exportación directa
        ot_master_csv = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid=22353124"
        procesos_csv = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid=1564553976"
        
        # Cargar datos
        ot_master = pd.read_csv(ot_master_csv)
        procesos = pd.read_csv(procesos_csv)
        
        return ot_master, procesos
    except Exception as e:
        st.error(f"Error al cargar los datos desde Google Sheets: {e}")
        return None, None

# Cargar datos con spinner
with st.spinner("Cargando datos desde Google Sheets..."):
    ot_master, procesos = load_data()

if ot_master is None or procesos is None:
    st.error("No se pudieron cargar los datos. Por favor, verifica la conexión e intenta nuevamente.")
    st.stop()

# Asegurar que la columna 'ot' sea string en ambos dataframes
ot_master['ot'] = ot_master['ot'].astype(str)
procesos['ot'] = procesos['ot'].astype(str)

# Sidebar con filtros
st.sidebar.header("🔍 Filtros")

# Convertir fechas en ot_master
date_columns = ['fecha_entrega', 'fecha_impresion', 'fecha_terminado', 'fecha_entregada']
for col in date_columns:
    if col in ot_master.columns:
        ot_master[col] = pd.to_datetime(ot_master[col], errors='coerce')

# Convertir fechas en procesos
date_columns_procesos = ['fecha_inicio_1', 'fecha_inicio_2']
for col in date_columns_procesos:
    if col in procesos.columns:
        procesos[col] = pd.to_datetime(procesos[col], errors='coerce')

# Filtros principales
clientes = ['Todos'] + sorted(ot_master['cliente'].dropna().unique().tolist())
cliente_seleccionado = st.sidebar.selectbox("Cliente", clientes)

estatus_options = ['Todos'] + sorted(ot_master['estatus'].dropna().unique().tolist())
estatus_seleccionado = st.sidebar.selectbox("Estatus", estatus_options)

# Filtro de OT
ots = ["Todas"] + sorted(ot_master['ot'].astype(str).unique().tolist())
ot_seleccionada = st.sidebar.selectbox("OT", ots)

# Filtros de empleados SIN REPETIDOS
st.sidebar.subheader("👥 Filtros por Empleados")

def limpiar_nombre(nombre):
    if pd.isna(nombre) or nombre == '' or nombre == ' ':
        return None
    nombre_limpio = str(nombre)
    nombre_limpio = nombre_limpio.strip()
    nombre_limpio = ' '.join(nombre_limpio.split())
    caracteres_problematicos = ['\n', '\t', '\r', '*', '#', '  ']
    for char in caracteres_problematicos:
        nombre_limpio = nombre_limpio.replace(char, ' ')
    nombre_limpio = nombre_limpio.title()
    nombre_limpio = ' '.join(nombre_limpio.split())
    return nombre_limpio if nombre_limpio != '' else None

# Obtener lista única de empleados
empleados_1 = [limpiar_nombre(x) for x in procesos['empleado_1'].dropna().unique()]
empleados_2 = [limpiar_nombre(x) for x in procesos['empleado_2'].dropna().unique()]
todos_empleados = list(set([emp for emp in empleados_1 + empleados_2 if emp is not None]))
todos_empleados = ['Todos'] + sorted(todos_empleados)

empleado_seleccionado = st.sidebar.selectbox("Empleado", todos_empleados)

# Filtro de fechas
st.sidebar.subheader("📅 Filtro por Fecha de Entrega")
min_date = ot_master['fecha_entrega'].min()
max_date = ot_master['fecha_entrega'].max()
if pd.notna(min_date) and pd.notna(max_date):
    fecha_inicio = st.sidebar.date_input("Fecha inicio", min_date)
    fecha_fin = st.sidebar.date_input("Fecha fin", max_date)
else:
    st.sidebar.warning("No hay fechas válidas para filtrar")
    fecha_inicio = None
    fecha_fin = None

# Aplicar filtros
ot_master_filtrado = ot_master.copy()
procesos_filtrados = procesos.copy()

if cliente_seleccionado != 'Todos':
    ot_master_filtrado = ot_master_filtrado[ot_master_filtrado['cliente'] == cliente_seleccionado]
    procesos_filtrados = procesos_filtrados[procesos_filtrados['ot'].isin(ot_master_filtrado['ot'])]

if estatus_seleccionado != 'Todos':
    ot_master_filtrado = ot_master_filtrado[ot_master_filtrado['estatus'] == estatus_seleccionado]
    procesos_filtrados = procesos_filtrados[procesos_filtrados['ot'].isin(ot_master_filtrado['ot'])]

if ot_seleccionada != 'Todas':
    ot_master_filtrado = ot_master_filtrado[ot_master_filtrado['ot'] == ot_seleccionada]
    procesos_filtrados = procesos_filtrados[procesos_filtrados['ot'] == ot_seleccionada]

if empleado_seleccionado != 'Todos':
    procesos_temp = procesos_filtrados.copy()
    procesos_temp['empleado_1_clean'] = procesos_temp['empleado_1'].apply(limpiar_nombre)
    procesos_temp['empleado_2_clean'] = procesos_temp['empleado_2'].apply(limpiar_nombre)
    procesos_filtrados = procesos_temp[
        (procesos_temp['empleado_1_clean'] == empleado_seleccionado) | 
        (procesos_temp['empleado_2_clean'] == empleado_seleccionado)
    ]
    procesos_filtrados = procesos_filtrados.drop(['empleado_1_clean', 'empleado_2_clean'], axis=1)
    ot_master_filtrado = ot_master_filtrado[ot_master_filtrado['ot'].isin(procesos_filtrados['ot'])]

if fecha_inicio and fecha_fin:
    ot_master_filtrado = ot_master_filtrado[
        (ot_master_filtrado['fecha_entrega'] >= pd.Timestamp(fecha_inicio)) &
        (ot_master_filtrado['fecha_entrega'] <= pd.Timestamp(fecha_fin))
    ]
    procesos_filtrados = procesos_filtrados[procesos_filtrados['ot'].isin(ot_master_filtrado['ot'])]

# Definir estados que NO se consideran vencidos
estados_no_vencidos = ['FACTURADO', 'OK', 'OK NO ENTREGADO']

# Calcular OTs vencidas y por vencer
hoy = datetime.now()
ot_master_filtrado['estado_entrega'] = ot_master_filtrado.apply(
    lambda row: 
        'Completada' if row['estatus'] in estados_no_vencidos else
        'Vencida' if pd.notna(row['fecha_entrega']) and row['fecha_entrega'] < hoy else 
        'Por vencer' if pd.notna(row['fecha_entrega']) and row['fecha_entrega'] >= hoy and row['fecha_entrega'] <= hoy + timedelta(days=7) else 
        'En plazo',
    axis=1
)

# Calcular porcentaje de facturación
total_ots = len(ot_master_filtrado)
ots_facturadas = len(ot_master_filtrado[ot_master_filtrado['estatus'] == 'FACTURADO'])
porcentaje_facturado = (ots_facturadas / total_ots * 100) if total_ots > 0 else 0

# Identificar reprocesos (Garantías)
if 'orden_compra' in ot_master_filtrado.columns:
    ot_master_filtrado['es_reproceso'] = ot_master_filtrado['orden_compra'].str.contains('GARANTIA', case=False, na=False)
    total_reprocesos = ot_master_filtrado['es_reproceso'].sum()
    porcentaje_reprocesos = (total_reprocesos / total_ots * 100) if total_ots > 0 else 0
else:
    total_reprocesos = 0
    porcentaje_reprocesos = 0

# Calcular desviaciones de horas
ots_desviacion_positiva = pd.DataFrame()
ots_desviacion_negativa = pd.DataFrame()

if 'horas_estimadas_ot' in ot_master_filtrado.columns and 'horas_reales_ot' in ot_master_filtrado.columns:
    # Filtrar solo OTs con horas válidas
    ot_con_horas = ot_master_filtrado[
        (ot_master_filtrado['horas_estimadas_ot'].notna()) & 
        (ot_master_filtrado['horas_reales_ot'].notna())
    ].copy()
    
    # Calcular desviaciones
    ot_con_horas['diferencia_horas'] = ot_con_horas['horas_reales_ot'] - ot_con_horas['horas_estimadas_ot']
    ot_con_horas['tipo_desviacion'] = ot_con_horas['diferencia_horas'].apply(
        lambda x: 'Desviación Positiva' if x <= 0 else 'Desviación Negativa'
    )
    
    # Separar en DataFrames para desviaciones positivas y negativas
    ots_desviacion_positiva = ot_con_horas[ot_con_horas['tipo_desviacion'] == 'Desviación Positiva'].copy()
    ots_desviacion_negativa = ot_con_horas[ot_con_horas['tipo_desviacion'] == 'Desviación Negativa'].copy()
    
    # Calcular totales
    total_horas_programadas = ot_con_horas['horas_estimadas_ot'].sum()
    horas_desviacion_positiva = ots_desviacion_positiva['horas_reales_ot'].sum()
    horas_desviacion_negativa = ots_desviacion_negativa['horas_reales_ot'].sum()
    
    # Calcular porcentajes
    porcentaje_positivo = (horas_desviacion_positiva / total_horas_programadas * 100) if total_horas_programadas > 0 else 0
    porcentaje_negativo = (horas_desviacion_negativa / total_horas_programadas * 100) if total_horas_programadas > 0 else 0
else:
    total_horas_programadas = 0
    horas_desviacion_positiva = 0
    horas_desviacion_negativa = 0
    porcentaje_positivo = 0
    porcentaje_negativo = 0

# Métricas principales
st.header("📊 Métricas Principales")
col1, col2, col3, col4, col5, col6 = st.columns(6)
with col1: 
    st.metric("Total OTs", total_ots)
with col2: 
    ots_en_proceso = len(ot_master_filtrado[ot_master_filtrado['estatus'] == 'EN PROCESO'])
    st.metric("OTs en Proceso", ots_en_proceso)
with col3:
    st.metric("OTs Facturadas", ots_facturadas, f"{porcentaje_facturado:.1f}%")
with col4: 
    ots_vencidas = len(ot_master_filtrado[(ot_master_filtrado['estado_entrega'] == 'Vencida') & (~ot_master_filtrado['estatus'].isin(estados_no_vencidos))])
    st.metric("OTs Vencidas", ots_vencidas, delta=-ots_vencidas, delta_color="inverse")
with col5: 
    ots_por_vencer = len(ot_master_filtrado[(ot_master_filtrado['estado_entrega'] == 'Por vencer') & (~ot_master_filtrado['estatus'].isin(estados_no_vencidos))])
    st.metric("OTs por Vencer", ots_por_vencer, delta=ots_por_vencer, delta_color="off")
with col6:
    st.metric("Reprocesos", total_reprocesos, f"{porcentaje_reprocesos:.1f}%")

# =============================================
# SECCIÓN DE EXPORTACIÓN (COLOCADA AQUÍ PARA MAYOR VISIBILIDAD)
# =============================================

st.markdown("---")
st.header("🚀 Exportar Reportes Ejecutivos")

# Verificación de funciones (DEBUG)
with st.expander("🔍 Verificación de Funciones (Debug)"):
    try:
        exportar_a_powerpoint
        st.success("✅ Función exportar_a_powerpoint: DISPONIBLE")
    except Exception as e:
        st.error(f"❌ Función exportar_a_powerpoint: NO DISPONIBLE - {e}")

    try:
        exportar_a_pdf
        st.success("✅ Función exportar_a_pdf: DISPONIBLE")
    except Exception as e:
        st.error(f"❌ Función exportar_a_pdf: NO DISPONIBLE - {e}")

    try:
        exportar_a_excel
        st.success("✅ Función exportar_a_excel: DISPONIBLE")
    except Exception as e:
        st.error(f"❌ Función exportar_a_excel: NO DISPONIBLE - {e}")

col1, col2, col3 = st.columns(3)

with col1:
    st.subheader("📊 PowerPoint Ejecutivo")
    st.info("Presentación profesional lista para reuniones")
    
    if st.button("🎯 Generar PowerPoint", use_container_width=True, type="primary", key="ppt_btn"):
        with st.spinner("Generando PowerPoint..."):
            try:
                exportar_a_powerpoint()
            except Exception as e:
                st.error(f"Error al generar PowerPoint: {e}")

with col2:
    st.subheader("📄 Reporte PDF")
    st.info("Documento formal para distribución")
    
    if st.button("📋 Generar PDF", use_container_width=True, key="pdf_btn"):
        with st.spinner("Generando PDF..."):
            try:
                exportar_a_pdf()
            except Exception as e:
                st.error(f"Error al generar PDF: {e}")

with col3:
    st.subheader("📈 Datos para Análisis")
    st.info("Datos completos en Excel para análisis detallado")
    
    if st.button("📊 Generar Reporte Excel", use_container_width=True, key="excel_btn"):
        with st.spinner("Generando Excel..."):
            try:
                exportar_a_excel()
            except Exception as e:
                st.error(f"Error al generar Excel: {e}")

# Información adicional
st.markdown("---")
st.info("""
💡 **Características de los Reportes:**

**PowerPoint Ejecutivo:**
- ✅ Presentación lista para reuniones
- ✅ Métricas clave resumidas
- ✅ Análisis de eficiencia
- ✅ OTs críticas identificadas
- ✅ Recomendaciones de acción

**Reporte PDF:**
- ✅ Documento formal listo para imprimir
- ✅ Métricas principales organizadas
- ✅ Análisis detallado
- ✅ Recomendaciones específicas

**Excel Completo:**
- ✅ Todos los datos filtrados
- ✅ Hojas organizadas por categoría
- ✅ Resumen ejecutivo incluido
- ✅ Formato listo para análisis
""")

st.markdown("---")

# CONTINÚA CON EL RESTO DEL DASHBOARD...

# GRÁFICO PRINCIPAL: OTs VENCIDAS Y POR VENCER
st.header("📅 Estado de Entregas - OTs Vencidas y Por Vencer")
estado_entrega_counts = ot_master_filtrado['estado_entrega'].value_counts()
estados_interes = ['Vencida', 'Por vencer']
estado_entrega_counts_filtrado = estado_entrega_counts[estado_entrega_counts.index.isin(estados_interes)]

if not estado_entrega_counts_filtrado.empty:
    fig_ots_vencidas = px.bar(
        x=estado_entrega_counts_filtrado.index,
        y=estado_entrega_counts_filtrado.values,
        title="OTs Vencidas y Por Vencer (Solo OTs Activas)",
        labels={'x': 'Estado de Entrega', 'y': 'Cantidad de OTs'},
        color=estado_entrega_counts_filtrado.index,
        color_discrete_map={'Vencida': '#FF4B4B', 'Por vencer': '#FFA500'},
        text=estado_entrega_counts_filtrado.values
    )
    fig_ots_vencidas.update_traces(texttemplate='%{text}', textposition='outside')
    fig_ots_vencidas.update_layout(showlegend=False, yaxis_title="Cantidad de OTs", xaxis_title="", height=400)
    st.plotly_chart(fig_ots_vencidas, use_container_width=True)
else:
    st.info("No hay OTs vencidas o por vencer con los filtros actuales.")

# ... (el resto del código permanece igual, incluyendo las funciones de exportación al final)

# =============================================
# FUNCIONES DE EXPORTACIÓN MEJORADAS
# =============================================

def exportar_a_powerpoint():
    """Exportar reporte ejecutivo a PowerPoint"""
    try:
        # Crear nueva presentación
        prs = Presentation()
        
        # Slide 1: Portada
        slide_layout = prs.slide_layouts[0]  # Layout de título
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Reporte de Producción"
        subtitle.text = f"Adimatec - {datetime.now().strftime('%d/%m/%Y')}"
        
        # Slide 2: Métricas Principales
        slide_layout = prs.slide_layouts[1]  # Layout de título y contenido
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Métricas Principales"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = f"""• Total OTs: {total_ots}
• OTs Facturadas: {ots_facturadas} ({porcentaje_facturado:.1f}%)
• OTs en Proceso: {ots_en_proceso}
• OTs Vencidas: {ots_vencidas}
• OTs por Vencer: {ots_por_vencer}
• Reprocesos: {total_reprocesos} ({porcentaje_reprocesos:.1f}%)"""
        
        # Slide 3: Análisis de Eficiencia
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Análisis de Eficiencia"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = f"""• Eficiencia de Facturación: {porcentaje_facturado:.1f}%
• Tasa de Reprocesos: {porcentaje_reprocesos:.1f}%
• Horas Programadas Totales: {total_horas_programadas:.1f}h
• Desviaciones Positivas: {porcentaje_positivo:.1f}%
• Desviaciones Negativas: {porcentaje_negativo:.1f}%"""
        
        # Slide 4: OTs Críticas (si existen)
        if not ots_desviacion_negativa.empty:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "OTs con Desviaciones Negativas"
            
            content = slide.placeholders[1]
            text_frame = content.text_frame
            
            # Tomar las 5 OTs con mayores desviaciones
            top_ots = ots_desviacion_negativa.nlargest(5, 'diferencia_horas')
            texto_ots = "Principales OTs con desviaciones:\n\n"
            for idx, row in top_ots.iterrows():
                texto_ots += f"• OT {row['ot']}: {row['diferencia_horas']:.1f}h (Cliente: {row.get('cliente', 'N/A')})\n"
            
            text_frame.text = texto_ots
        
        # Slide 5: Recomendaciones
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Recomendaciones y Acciones"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = """• Enfocar recursos en OTs vencidas y por vencer
• Analizar causas de reprocesos
• Optimizar estimación de horas
• Revisar OTs con mayores desviaciones
• Mantener comunicación con clientes críticos"""
        
        # Guardar en memoria
        from io import BytesIO
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Botón de descarga
        st.download_button(
            label="🎯 Descargar PowerPoint Ejecutivo",
            data=pptx_buffer.getvalue(),
            file_name=f"Reporte_Ejecutivo_Adimatec_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        
        st.success("✅ Presentación PowerPoint generada exitosamente!")
        
    except Exception as e:
        st.error(f"Error al generar PowerPoint: {str(e)}")

def exportar_a_pdf():
    """Exportar reporte a PDF"""
    try:
        # Crear PDF
        pdf = FPDF()
        pdf.add_page()
        
        # Encabezado
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(200, 10, "Reporte de Producción - Adimatec", 0, 1, 'C')
        pdf.ln(5)
        
        pdf.set_font("Arial", '', 12)
        pdf.cell(200, 10, f"Generado el: {datetime.now().strftime('%d/%m/%Y %H:%M')}", 0, 1, 'C')
        pdf.ln(10)
        
        # Métricas Principales
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(200, 10, "Métricas Principales", 0, 1)
        pdf.ln(5)
        
        pdf.set_font("Arial", '', 12)
        metricas = [
            f"Total OTs: {total_ots}",
            f"OTs Facturadas: {ots_facturadas} ({porcentaje_facturado:.1f}%)",
            f"OTs en Proceso: {ots_en_proceso}",
            f"OTs Vencidas: {ots_vencidas}",
            f"OTs por Vencer: {ots_por_vencer}",
            f"Reprocesos: {total_reprocesos} ({porcentaje_reprocesos:.1f}%)"
        ]
        
        for metrica in metricas:
            pdf.cell(200, 10, metrica, 0, 1)
        
        pdf.ln(10)
        
        # Análisis de Eficiencia
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(200, 10, "Análisis de Eficiencia", 0, 1)
        pdf.ln(5)
        
        pdf.set_font("Arial", '', 12)
        eficiencia = [
            f"Horas Programadas Totales: {total_horas_programadas:.1f}h",
            f"Desviaciones Positivas: {porcentaje_positivo:.1f}%",
            f"Desviaciones Negativas: {porcentaje_negativo:.1f}%"
        ]
        
        for item in eficiencia:
            pdf.cell(200, 10, item, 0, 1)
        
        pdf.ln(10)
        
        # Recomendaciones
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(200, 10, "Recomendaciones", 0, 1)
        pdf.ln(5)
        
        pdf.set_font("Arial", '', 12)
        recomendaciones = [
            "• Enfocar recursos en OTs vencidas y por vencer",
            "• Analizar causas de reprocesos",
            "• Optimizar estimación de horas",
            "• Revisar OTs con mayores desviaciones",
            "• Mantener comunicación con clientes críticos"
        ]
        
        for recomendacion in recomendaciones:
            pdf.cell(200, 10, recomendacion, 0, 1)
        
        # Guardar en memoria
        pdf_output = "reporte_adimatec.pdf"
        pdf.output(pdf_output)
        
        with open(pdf_output, "rb") as f:
            st.download_button(
                label="📄 Descargar PDF",
                data=f.read(),
                file_name=f"Reporte_Adimatec_{datetime.now().strftime('%Y%m%d')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )
        
        # Limpiar archivo temporal
        if os.path.exists(pdf_output):
            os.remove(pdf_output)
            
        st.success("✅ Reporte PDF generado exitosamente!")
        
    except Exception as e:
        st.error(f"Error al generar PDF: {str(e)}")

def exportar_a_excel():
    """Exportar datos completos a Excel"""
    try:
        # Crear un escritor de Excel
        with pd.ExcelWriter('reporte_adimatec.xlsx', engine='openpyxl') as writer:
            # Hoja 1: OT Master
            ot_master_filtrado.to_excel(writer, sheet_name='OT_Master', index=False)
            
            # Hoja 2: Procesos
            if not procesos_filtrados.empty:
                procesos_filtrados.to_excel(writer, sheet_name='Procesos', index=False)
            
            # Hoja 3: Resumen Ejecutivo
            resumen_data = {
                'Métrica': [
                    'Total OTs', 
                    'OTs Facturadas', 
                    'OTs en Proceso', 
                    'OTs Vencidas', 
                    'OTs por Vencer',
                    '% Facturación',
                    '% Reprocesos',
                    'Horas Programadas Totales',
                    'Desviaciones Positivas',
                    'Desviaciones Negativas'
                ],
                'Valor': [
                    total_ots,
                    ots_facturadas,
                    ots_en_proceso,
                    ots_vencidas,
                    ots_por_vencer,
                    f"{porcentaje_facturado:.1f}%",
                    f"{porcentaje_reprocesos:.1f}%",
                    f"{total_horas_programadas:.1f}h",
                    f"{horas_desviacion_positiva:.1f}h",
                    f"{horas_desviacion_negativa:.1f}h"
                ]
            }
            pd.DataFrame(resumen_data).to_excel(writer, sheet_name='Resumen', index=False)
            
            # Hoja 4: OTs Críticas
            if not ots_desviacion_negativa.empty:
                columnas_criticas = ['ot', 'cliente', 'horas_estimadas_ot', 'horas_reales_ot', 'diferencia_horas']
                columnas_disponibles = [col for col in columnas_criticas if col in ots_desviacion_negativa.columns]
                if columnas_disponibles:
                    ots_desviacion_negativa[columnas_disponibles].to_excel(writer, sheet_name='OTs_Criticas', index=False)
        
        # Ofrecer descarga
        with open('reporte_adimatec.xlsx', 'rb') as f:
            st.download_button(
                label="📈 Descargar Excel Completo",
                data=f.read(),
                file_name=f"Reporte_Adimatec_{datetime.now().strftime('%Y%m%d')}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
        
        # Limpiar archivo temporal
        if os.path.exists('reporte_adimatec.xlsx'):
            os.remove('reporte_adimatec.xlsx')
            
        st.success("✅ Archivo Excel generado exitosamente!")
        
    except Exception as e:
        st.error(f"Error al generar Excel: {str(e)}")

# ... (el resto del código, tablas de datos, footer, etc.)

# Tablas de datos
st.markdown("---")
st.header("📋 Datos Detallados")
tab1, tab2 = st.tabs(["OT Master", "Procesos"])
with tab1:
    st.subheader("Tabla OT Master")
    columnas_mostrar = ['ot', 'descripcion', 'cliente', 'estatus', 'fecha_entrega', 'horas_estimadas_ot', 'horas_reales_ot']
    columnas_disponibles = [col for col in columnas_mostrar if col in ot_master_filtrado.columns]
    if not ot_master_filtrado.empty:
        st.dataframe(ot_master_filtrado[columnas_disponibles], use_container_width=True, hide_index=True)
        csv_ot = ot_master_filtrado.to_csv(index=False)
        st.download_button(label="📥 Descargar OT Master como CSV", data=csv_ot, file_name="ot_master_filtrado.csv", mime="text/csv")
    else: 
        st.info("No hay datos para mostrar en OT Master")
with tab2:
    st.subheader("Tabla Procesos")
    posibles_nombres = ['proceso', 'Proceso', 'PROCESO', 'proceso_nombre', 'Proceso_Nombre']
    columna_proceso = None
    for nombre in posibles_nombres:
        if nombre in procesos_filtrados.columns:
            columna_proceso = nombre
            break
    columnas_mostrar_procesos = ['ot', columna_proceso, 'horas_estimadas', 'horas_reales', 'empleado_1', 'empleado_2']
    columnas_disponibles_procesos = [col for col in columnas_mostrar_procesos if col in procesos_filtrados.columns]
    if not procesos_filtrados.empty:
        st.dataframe(procesos_filtrados[columnas_disponibles_procesos], use_container_width=True, hide_index=True)
        csv_procesos = procesos_filtrados.to_csv(index=False)
        st.download_button(label="📥 Descargar Procesos como CSV", data=csv_procesos, file_name="procesos_filtrados.csv", mime="text/csv")
    else: 
        st.info("No hay datos para mostrar en Procesos")

# Footer
st.markdown("---")
st.markdown(
    """
    <div style='text-align: center'>
        <p>Dashboard de Producción - Adimatec | Desarrollado con Streamlit</p>
        <p><small>✨ Incluye análisis de Pareto y exportación a PowerPoint/PDF</small></p>
    </div>
    """,
    unsafe_allow_html=True
)
